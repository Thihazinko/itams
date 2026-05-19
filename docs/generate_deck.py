"""
Generate a modern 17-slide PowerPoint deck for Infra Ninja / ITAMS.

Run from project root:
    python docs/generate_deck.py

Output: docs/InfraNinja_ITAMS_Overview.pptx
"""
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree

SHOTS = Path(__file__).resolve().parent / "screenshots"
TOTAL = 22

# ---------- typography ----------
FONT = "Segoe UI"

# ---------- color palette ----------
COLOR_PRIMARY = RGBColor(0x0D, 0x6E, 0xFD)   # blue
COLOR_ACCENT  = RGBColor(0x6F, 0x3C, 0xFF)   # purple
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)   # green
COLOR_WARN    = RGBColor(0xF5, 0x9E, 0x0B)   # amber
COLOR_DANGER  = RGBColor(0xEF, 0x44, 0x44)   # red
COLOR_TEXT    = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
COLOR_TEXT_2  = RGBColor(0x33, 0x40, 0x55)   # slate-700
COLOR_MUTED   = RGBColor(0x64, 0x74, 0x8B)   # slate-500
COLOR_FAINT   = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
COLOR_BG_TOP  = RGBColor(0xFA, 0xFB, 0xFC)   # near-white
COLOR_BG_BOT  = RGBColor(0xEE, 0xF1, 0xF6)   # slightly cooler
COLOR_HAIRLINE= RGBColor(0xE2, 0xE8, 0xF0)
COLOR_DARK    = RGBColor(0x0F, 0x17, 0x2A)
COLOR_DARK_2  = RGBColor(0x1A, 0x21, 0x30)
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# section pill: (label, color)  per slide number (1-indexed)
SECTION = {
    1:  ("INTRO",        COLOR_ACCENT),
    2:  ("OVERVIEW",     COLOR_ACCENT),
    3:  ("OVERVIEW",     COLOR_ACCENT),   # Team
    4:  ("OVERVIEW",     COLOR_ACCENT),   # Tech stack
    5:  ("OVERVIEW",     COLOR_ACCENT),   # Architecture
    6:  ("MODULES",      COLOR_SUCCESS),  # ── DIVIDER ──
    7:  ("MODULES",      COLOR_SUCCESS),
    8:  ("MODULES",      COLOR_SUCCESS),
    9:  ("MODULES",      COLOR_SUCCESS),
    10: ("USERS & AUTH", COLOR_PRIMARY),  # ── DIVIDER ──
    11: ("USERS & AUTH", COLOR_PRIMARY),
    12: ("USERS & AUTH", COLOR_PRIMARY),
    13: ("USERS & AUTH", COLOR_PRIMARY),
    14: ("INSIGHTS",     COLOR_WARN),     # ── DIVIDER ──
    15: ("INSIGHTS",     COLOR_WARN),
    16: ("INSIGHTS",     COLOR_WARN),
    17: ("OPERATIONS",   COLOR_DANGER),   # ── DIVIDER ──
    18: ("OPERATIONS",   COLOR_DANGER),
    19: ("OPERATIONS",   COLOR_DANGER),
    20: ("WRAP-UP",      COLOR_MUTED),    # ── DIVIDER ──
    21: ("WRAP-UP",      COLOR_MUTED),
    22: ("WRAP-UP",      COLOR_MUTED),
}


# ============================================================
# Low-level XML helpers (for gradients, shadows, picture fills)
# ============================================================
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def _a(tag):
    return f"{{{NSMAP['a']}}}{tag}"


def set_gradient_fill(shape, color_top, color_bottom, angle_deg=270):
    """Replace the shape's fill with a top-to-bottom linear gradient."""
    sppr = shape.fill._xPr.find(_a('solidFill'))
    if sppr is not None:
        shape.fill._xPr.remove(sppr)
    grad = shape.fill._xPr.find(_a('gradFill'))
    if grad is not None:
        shape.fill._xPr.remove(grad)
    grad = etree.SubElement(shape.fill._xPr, _a('gradFill'))
    grad.set('flip', 'none')
    grad.set('rotWithShape', '1')
    stops = etree.SubElement(grad, _a('gsLst'))
    for pos, col in ((0, color_top), (100000, color_bottom)):
        gs = etree.SubElement(stops, _a('gs'))
        gs.set('pos', str(pos))
        srgb = etree.SubElement(gs, _a('srgbClr'))
        srgb.set('val', f'{col[0]:02X}{col[1]:02X}{col[2]:02X}')
    lin = etree.SubElement(grad, _a('lin'))
    lin.set('ang', str(angle_deg * 60000))
    lin.set('scaled', '1')
    tile = etree.SubElement(grad, _a('tileRect'))


def add_soft_shadow(shape, blur_pt=14, dist_pt=6, color=COLOR_TEXT, alpha=22):
    """Append an outer shadow to a shape's <a:spPr>. Subtle drop shadow."""
    spPr = shape._element.find('.//' + _a('spPr'))
    if spPr is None:
        return
    # remove existing effectLst if present
    existing = spPr.find(_a('effectLst'))
    if existing is not None:
        spPr.remove(existing)
    effect = etree.SubElement(spPr, _a('effectLst'))
    outer = etree.SubElement(effect, _a('outerShdw'))
    outer.set('blurRad', str(blur_pt * 12700))
    outer.set('dist', str(dist_pt * 12700))
    outer.set('dir', '5400000')   # straight down
    outer.set('rotWithShape', '0'); outer.set('algn', 'ctr')
    srgb = etree.SubElement(outer, _a('srgbClr'))
    srgb.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')
    alpha_el = etree.SubElement(srgb, _a('alpha'))
    alpha_el.set('val', str(alpha * 1000))


def set_slide_bg_solid(slide, rgb):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb


def set_slide_bg_gradient(slide, top, bottom):
    """Slide background gradient using OXML."""
    bg = slide.background
    fill = bg.fill
    fill.solid()  # ensure xPr is initialized
    # use the low-level helper on the fill
    fill_xpr = fill._xPr
    # remove solidFill
    sf = fill_xpr.find(_a('solidFill'))
    if sf is not None:
        fill_xpr.remove(sf)
    # build gradient
    grad = etree.SubElement(fill_xpr, _a('gradFill'))
    grad.set('flip', 'none')
    grad.set('rotWithShape', '1')
    stops = etree.SubElement(grad, _a('gsLst'))
    for pos, col in ((0, top), (100000, bottom)):
        gs = etree.SubElement(stops, _a('gs'))
        gs.set('pos', str(pos))
        srgb = etree.SubElement(gs, _a('srgbClr'))
        srgb.set('val', f'{col[0]:02X}{col[1]:02X}{col[2]:02X}')
    lin = etree.SubElement(grad, _a('lin'))
    lin.set('ang', str(270 * 60000))
    lin.set('scaled', '1')
    etree.SubElement(grad, _a('tileRect'))


def style_run(run, *, size=12, bold=False, italic=False, color=None, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    if bold: run.font.bold = True
    if italic: run.font.italic = True
    if color is not None: run.font.color.rgb = color


# ============================================================
# Template chrome: pill, title, subtitle, footer
# ============================================================
def add_decorative_blob(slide, prs, color, x_offset_in=8.5, y_offset_in=-2.5, size_in=5.5):
    """Soft colored circle bleeding off a corner — purely decorative."""
    blob = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x_offset_in), Inches(y_offset_in),
                                  Inches(size_in), Inches(size_in))
    blob.fill.solid()
    blob.fill.fore_color.rgb = color
    blob.line.fill.background()
    # subtle transparency via alpha
    sp = blob.fill._xPr.find(_a('solidFill'))
    if sp is not None:
        clr = sp.find(_a('srgbClr'))
        if clr is not None:
            alpha = etree.SubElement(clr, _a('alpha'))
            alpha.set('val', '12000')   # ~12% opacity
    return blob


def add_section_pill(slide, prs, label, page_num, color):
    """Small uppercase pill in top-left: '06 · MODULES'."""
    pill_text = f"{page_num:02d}  ·  {label}"
    # measure-ish: each char ~6px at 9pt, plus padding
    w_in = max(1.4, 0.06 * len(pill_text) + 0.5)
    x, y = Inches(0.55), Inches(0.4)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, y, Inches(w_in), Inches(0.32))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    # soften via alpha
    sp = pill.fill._xPr.find(_a('solidFill'))
    if sp is not None:
        clr = sp.find(_a('srgbClr'))
        if clr is not None:
            alpha = etree.SubElement(clr, _a('alpha'))
            alpha.set('val', '15000')   # 15%
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = pill_text
    style_run(r, size=9, bold=True, color=color)


def add_modern_title(slide, text, top=0.95, size=32):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.85))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    style_run(r, size=size, bold=True, color=COLOR_TEXT)
    return box


def add_modern_subtitle(slide, text, top=1.65, size=13):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.5))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    style_run(r, size=size, color=COLOR_MUTED)
    return box


def add_progress_bar(slide, prs, current, total=TOTAL, color=COLOR_PRIMARY):
    """Bottom-right segmented progress bar (small dots) — modern slide indicator."""
    dot_size = Inches(0.085)
    gap = Inches(0.04)
    total_w = dot_size * total + gap * (total - 1)
    x0 = prs.slide_width - Inches(0.55) - total_w
    y = prs.slide_height - Inches(0.35)
    for i in range(total):
        x = x0 + i * (dot_size + gap)
        is_active = (i + 1) == current
        is_past = (i + 1) < current
        seg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, dot_size, dot_size)
        seg.fill.solid()
        if is_active:
            seg.fill.fore_color.rgb = color
        elif is_past:
            seg.fill.fore_color.rgb = COLOR_FAINT
        else:
            seg.fill.fore_color.rgb = COLOR_HAIRLINE
        seg.line.fill.background()


def add_footer_modern(slide, prs, page_num, total=TOTAL, color=COLOR_PRIMARY):
    """Project name (bottom-left), progress dots (bottom-right)."""
    # left text
    left = slide.shapes.add_textbox(Inches(0.55),
                                    prs.slide_height - Inches(0.42),
                                    Inches(7), Inches(0.32))
    tf = left.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "INFRA NINJA  ·  ITAMS"
    style_run(r, size=8, bold=True, color=COLOR_FAINT)
    r.font._rPr.set('spc', '160')   # letter-spacing
    add_progress_bar(slide, prs, page_num, total, color)


def apply_template_chrome(slide, prs, page_num):
    """All the per-slide template furniture: bg gradient + pill + footer + decorative blob."""
    label, color = SECTION.get(page_num, ("", COLOR_PRIMARY))
    set_slide_bg_gradient(slide, COLOR_BG_TOP, COLOR_BG_BOT)
    add_decorative_blob(slide, prs, color)
    add_section_pill(slide, prs, label, page_num, color)
    add_footer_modern(slide, prs, page_num, color=color)


# ============================================================
# Content widgets
# ============================================================
def add_bullets(slide, items, left=Inches(0.55), top=Inches(2.3), width=Inches(11.8),
                height=Inches(4.5), size=14, color=COLOR_TEXT_2, spacing=8,
                bullet_color=None, two_tone_label_color=COLOR_TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            label, sub = item
            r = p.add_run(); r.text = "●  "
            style_run(r, size=size - 2, color=bullet_color or COLOR_PRIMARY, bold=True)
            r2 = p.add_run(); r2.text = label
            style_run(r2, size=size, bold=True, color=two_tone_label_color)
            if sub:
                r3 = p.add_run(); r3.text = "   " + sub
                style_run(r3, size=size - 1, color=color)
        else:
            r = p.add_run(); r.text = "●  "
            style_run(r, size=size - 2, color=bullet_color or COLOR_PRIMARY, bold=True)
            r2 = p.add_run(); r2.text = item
            style_run(r2, size=size, color=color)
        p.space_after = Pt(spacing)
    return box


def add_flat_card(slide, x, y, w, h, title, body_lines, accent=COLOR_PRIMARY):
    """Card with a left accent stripe, no top stripe. Flatter, more whitespace."""
    # main card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_HAIRLINE
    card.line.width = Pt(0.75)
    add_soft_shadow(card, blur_pt=18, dist_pt=4, alpha=10)
    # left accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.085), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    # content
    pad = Inches(0.28)
    inner = slide.shapes.add_textbox(x + Inches(0.32), y + pad,
                                     w - Inches(0.5), h - pad - Inches(0.15))
    tf = inner.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    style_run(r, size=14, bold=True, color=accent)
    p.space_after = Pt(6)
    for line in body_lines:
        bp = tf.add_paragraph()
        r1 = bp.add_run(); r1.text = "●  "
        style_run(r1, size=9, bold=True, color=accent)
        r2 = bp.add_run(); r2.text = line
        style_run(r2, size=11, color=COLOR_TEXT_2)
        bp.space_after = Pt(3)
    return card


def add_screenshot_rounded(slide, name, x, y, w, h):
    """Embed a PNG with rounded corners + subtle shadow.
    Works by adding a normal picture, then mutating its prstGeom from rect → roundRect."""
    path = SHOTS / name
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    # transform geometry to rounded-rect with a small corner radius
    prstGeom = pic._element.find('.//' + _a('prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        # ensure avLst exists; set adj to ~5% of the smaller dimension
        avLst = prstGeom.find(_a('avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, _a('avLst'))
        for old_gd in list(avLst.findall(_a('gd'))):
            avLst.remove(old_gd)
        gd = etree.SubElement(avLst, _a('gd'))
        gd.set('name', 'adj'); gd.set('fmla', 'val 5000')
    # subtle hairline border
    pic.line.color.rgb = COLOR_HAIRLINE
    pic.line.width = Pt(0.75)
    add_soft_shadow(pic, blur_pt=20, dist_pt=8, alpha=18)
    return pic


def add_caption(slide, text, x, y, w, color=COLOR_MUTED, size=10, align=PP_ALIGN.CENTER, italic=True):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.3))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    style_run(r, size=size, color=color, italic=italic)
    return box


def add_pill(slide, x, y, w, h, text, fill_color, text_color=COLOR_WHITE, size=11, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = fill_color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    style_run(r, size=size, bold=bold, color=text_color)
    return pill


# ============================================================
# Hero (dark) slide helpers
# ============================================================
HERO_BG_TOP    = RGBColor(0x33, 0x4D, 0xCC)   # deeper royal blue
HERO_BG_BOTTOM = RGBColor(0x4C, 0x1D, 0x95)   # deeper indigo / purple


def _set_blob_alpha(blob, alpha_pct=15):
    sp = blob.fill._xPr.find(_a('solidFill'))
    if sp is not None:
        clr = sp.find(_a('srgbClr'))
        if clr is not None:
            existing = clr.find(_a('alpha'))
            if existing is not None:
                clr.remove(existing)
            alpha = etree.SubElement(clr, _a('alpha'))
            alpha.set('val', str(alpha_pct * 1000))


def hero_slide(prs):
    """Brand gradient (deeper blue → indigo) with small accent blobs in the corners
    only — kept well clear of the central text area so white text stays high-contrast."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_gradient(s, HERO_BG_TOP, HERO_BG_BOTTOM)
    # Small accent in extreme top-right corner only (decorative, doesn't touch text)
    blob = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              prs.slide_width - Inches(2.5),
                              Inches(-2.2),
                              Inches(4.5), Inches(4.5))
    blob.fill.solid(); blob.fill.fore_color.rgb = RGBColor(0xA8, 0x7F, 0xFF)
    blob.line.fill.background()
    _set_blob_alpha(blob, alpha_pct=22)
    # Small accent in extreme bottom-left
    blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(-1.5),
                               prs.slide_height - Inches(2.5),
                               Inches(4), Inches(4))
    blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(0x6B, 0x8A, 0xFF)
    blob2.line.fill.background()
    _set_blob_alpha(blob2, alpha_pct=22)
    return s


# ============================================================
# Chapter divider slide
# ============================================================
def _lighten(color, amount=0.88):
    """Mix `color` with white to produce a faint tint (0 = original, 1 = white)."""
    return RGBColor(
        int(color[0] + (255 - color[0]) * amount),
        int(color[1] + (255 - color[1]) * amount),
        int(color[2] + (255 - color[2]) * amount),
    )


def slide_divider(prs, page_num, chapter_num, chapter_label, description, color):
    """Magazine-style chapter cover. Giant background numeral, big section title,
    one-line description. Subtle section-color decoration."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_gradient(s, COLOR_BG_TOP, COLOR_BG_BOT)

    # Giant background numeral (faint, magazine-watermark feel)
    num_box = s.shapes.add_textbox(Inches(0.4), Inches(0.2),
                                   Inches(7), Inches(5.0))
    tf = num_box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{chapter_num:02d}"
    style_run(r, size=320, bold=True, color=_lighten(color, 0.82))

    # Decorative blob bottom-right
    add_decorative_blob(s, prs, color, x_offset_in=8.5, y_offset_in=3, size_in=6)

    # "CHAPTER N" small pill, mid-left
    add_pill(s, Inches(0.65), Inches(3.65),
             Inches(1.7), Inches(0.34),
             f"CHAPTER  {chapter_num:02d}", color, size=9)

    # Big section title
    title = s.shapes.add_textbox(Inches(0.55), Inches(4.05),
                                 Inches(12.2), Inches(1.4))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = chapter_label
    style_run(r, size=64, bold=True, color=COLOR_TEXT)

    # Short accent underline
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(5.45),
                              Inches(1.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

    # One-line description
    desc = s.shapes.add_textbox(Inches(0.55), Inches(5.65),
                                Inches(11), Inches(0.7))
    p = desc.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = description
    style_run(r, size=15, color=COLOR_MUTED)

    # standard footer with progress
    add_footer_modern(s, prs, page_num, color=color)


# ============================================================
# Slides
# ============================================================
def slide1_title(prs):
    s = hero_slide(prs)

    # tiny eyebrow — saturated golden cream for high contrast
    eyebrow = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(12), Inches(0.4))
    p = eyebrow.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "IT ASSETS MANAGEMENT  ·  LARAVEL 11"
    style_run(r, size=11, bold=True, color=RGBColor(0xFC, 0xD3, 0x4D))
    r.font._rPr.set('spc', '320')

    # massive title — pure white
    title = s.shapes.add_textbox(Inches(0.7), Inches(2.85), Inches(12), Inches(1.7))
    p = title.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Infra Ninja"
    style_run(r, size=80, bold=True, color=COLOR_WHITE)

    # subtitle — pure white for max readability
    sub = s.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(12), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "PCs · Devices · Subscriptions · Licenses & Contracts"
    style_run(r, size=22, color=COLOR_WHITE)

    # tagline — soft but still legible
    tag = s.shapes.add_textbox(Inches(0.7), Inches(5.25), Inches(12), Inches(0.5))
    p = tag.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "One place. One workflow. Renewals that don't slip through the cracks."
    style_run(r, size=14, color=RGBColor(0xE0, 0xE7, 0xFF), italic=True)

    # bottom-right page indicator — high-contrast on the brand gradient
    page = s.shapes.add_textbox(prs.slide_width - Inches(1.0),
                                prs.slide_height - Inches(0.5),
                                Inches(0.6), Inches(0.3))
    p = page.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"01 / {TOTAL}"
    style_run(r, size=10, color=COLOR_WHITE, bold=True)


def slide2_overview(prs, n=2):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "What is Infra Ninja?")
    add_modern_subtitle(s, "A unified web app for tracking every renewable IT asset in your organization.")

    add_bullets(s, [
        ("Single source of truth",     "for PCs, peripherals, subscriptions and licenses"),
        ("Granular permissions",       "per-module View / Edit on top of admin / user roles"),
        ("Proactive renewals",         "live expiring badge, staggered email digests, mark-as-read"),
        ("Audit-ready",                "every change is logged — who, when, what"),
        ("Self-service",               "users update their own profile and password"),
        ("Modern UI",                  "Bootstrap 5 + glass-card auth + dark mode"),
    ], top=Inches(2.4), size=15, spacing=10)


def slide3_stack(prs, n=4):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Tech stack")
    add_modern_subtitle(s, "Boring, stable, well-supported. No exotic dependencies.")

    cards = [
        ("Backend",      ["PHP 8.2+", "Laravel 11.31", "Eloquent ORM"],                                COLOR_PRIMARY),
        ("Frontend",     ["Blade templates", "Bootstrap 5", "Bootstrap Icons", "Vite + Tailwind"],     COLOR_ACCENT),
        ("Database",     ["MySQL / MariaDB", "JSON columns", "Encrypted fields", "11 tables"],         COLOR_SUCCESS),
        ("Integrations", ["SMTP mail", "Excel I/O", "Chart.js", "Windows Task Scheduler"],             COLOR_WARN),
    ]
    x0, y0 = Inches(0.55), Inches(2.5)
    cw, ch = Inches(3.0), Inches(2.85)
    gap = Inches(0.15)
    for i, (title, body, c) in enumerate(cards):
        add_flat_card(s, x0 + i * (cw + gap), y0, cw, ch, title, body, accent=c)


def slide4_arch(prs, n=5):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "High-level architecture")
    add_modern_subtitle(s, "Classic Laravel MVC — session cookies, server-rendered Blade, no SPA.")

    # flow pills
    pills = [
        ("Browser  ·  Blade UI",   COLOR_TEXT,    Inches(0.7)),
        ("Laravel 11  ·  Routes / Controllers",   COLOR_PRIMARY, Inches(4.1)),
        ("MySQL  ·  Eloquent",   COLOR_SUCCESS, Inches(8.7)),
    ]
    widths = [Inches(2.85), Inches(4.2), Inches(2.85)]
    for (text, color, x), w in zip(pills, widths):
        add_pill(s, x, Inches(2.6), w, Inches(0.6), text, color, size=12)

    # arrows
    for x in [Inches(3.6), Inches(8.35)]:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(2.74), Inches(0.45), Inches(0.32))
        arr.fill.solid(); arr.fill.fore_color.rgb = COLOR_FAINT
        arr.line.fill.background()

    # cross-cutting concerns
    add_flat_card(s, Inches(0.55), Inches(3.7), Inches(6.0), Inches(2.5),
                  "Cross-cutting concerns",
                  ["Middleware: admin · module:<name>,<view|edit>",
                   "AppServiceProvider applies DB-driven mail config at boot",
                   "ActivityLogger captures every state change"],
                  accent=COLOR_ACCENT)
    add_flat_card(s, Inches(6.75), Inches(3.7), Inches(6.0), Inches(2.5),
                  "Scheduled work",
                  ["app:check-expirations runs once daily",
                   "Marks overdue Subscriptions",
                   "Sends staggered reminder digests (per day-mark)"],
                  accent=COLOR_DANGER)


def slide5_modules(prs, n=7):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Four asset modules")
    add_modern_subtitle(s, "Each has full CRUD, search, filter, Excel import / export, and bulk delete.")

    mods = [
        ("PC Master",           "Workstations, employees, status, encrypted credentials.", COLOR_PRIMARY, "/pc-assets"),
        ("Device Master",       "Peripherals, inventory qty, serial, warranty.",            COLOR_SUCCESS, "/devices"),
        ("Subscriptions",       "Domain, SSL, hosting, cloud, renewal cadence.",            COLOR_ACCENT,  "/subscriptions"),
        ("License & Contract",  "Software licenses & vendor contracts.",                    COLOR_WARN,    "/licenses-contracts"),
    ]
    x0, y0 = Inches(0.55), Inches(2.4)
    cw, ch = Inches(6.05), Inches(2.05)
    gap_x, gap_y = Inches(0.15), Inches(0.18)
    for i, (name, desc, color, path) in enumerate(mods):
        row, col = divmod(i, 2)
        x = x0 + col * (cw + gap_x)
        y = y0 + row * (ch + gap_y)
        add_flat_card(s, x, y, cw, ch, name, [desc, f"Route prefix: {path}"], accent=color)


def slide6_pc_devices(prs, n=8):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "PC Master & Device Master")
    add_modern_subtitle(s, "Hardware inventory — who has what, where it lives, when it was bought.")

    img_w = Inches(6.05); img_h = Inches(3.78)
    add_screenshot_rounded(s, "11-pc-assets.png",  Inches(0.55), Inches(2.4), img_w, img_h)
    add_caption(s, "PC Master — workstation list with status, OS, warranty",
                Inches(0.55), Inches(6.25), img_w)
    add_screenshot_rounded(s, "12-devices.png",     Inches(6.75), Inches(2.4), img_w, img_h)
    add_caption(s, "Device Master — peripherals with qty, vendor, delivery date",
                Inches(6.75), Inches(6.25), img_w)


def slide7_subs_lic(prs, n=9):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Subscriptions & License / Contract")
    add_modern_subtitle(s, "Recurring expenses that need renewing — the heart of the reminder system.")

    img_w = Inches(6.05); img_h = Inches(3.78)
    add_screenshot_rounded(s, "13-subscriptions.png", Inches(0.55), Inches(2.4), img_w, img_h)
    add_caption(s, "Subscriptions — domains, SSL, hosting, cloud", Inches(0.55), Inches(6.25), img_w)
    add_screenshot_rounded(s, "14-licenses-contracts.png", Inches(6.75), Inches(2.4), img_w, img_h)
    add_caption(s, "License & Contract — software licenses & vendor contracts", Inches(6.75), Inches(6.25), img_w)


def slide8_perms(prs, n=11):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Users & permissions")
    add_modern_subtitle(s, "Two roles + 8 granular flags = no shared admin password, no all-or-nothing.")

    # role pills row
    add_pill(s, Inches(0.55), Inches(2.45), Inches(2.6), Inches(0.6), "ADMIN  ·  full access", COLOR_WARN, size=12)
    add_pill(s, Inches(3.3),  Inches(2.45), Inches(3.4), Inches(0.6), "USER  ·  granted per module", COLOR_PRIMARY, size=12)

    add_bullets(s, [
        ("Per-module View",       "see index, show, export, template"),
        ("Per-module Edit",       "additionally allowed to create / update / delete / import"),
        ("Edit implies View",     "enforced in form, controller, and model"),
        ("Admin bypass",          "canView() / canEdit() short-circuit for admins"),
        ("Route middleware",      "module:<name>,<view|edit> wraps every module endpoint"),
        ("User management",       "admin-only at /users"),
    ], top=Inches(3.4), size=14, spacing=8)


def slide9_auth(prs, n=12):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Authentication flow")
    add_modern_subtitle(s, "Login + admin-only password reset by email + role-aware contact-admin branch.")

    img_w = Inches(4.0); img_h = Inches(2.5)
    y = Inches(2.4)
    xs = [Inches(0.55), Inches(4.7), Inches(8.85)]
    add_screenshot_rounded(s, "01-login.png",            xs[0], y, img_w, img_h)
    add_screenshot_rounded(s, "02-forgot-password.png",  xs[1], y, img_w, img_h)
    add_screenshot_rounded(s, "03-reset-password.png",   xs[2], y, img_w, img_h)
    add_caption(s, "1 ·  Sign in",         xs[0], Inches(5.0), img_w, color=COLOR_PRIMARY, size=12, italic=False)
    add_caption(s, "2 ·  Forgot password", xs[1], Inches(5.0), img_w, color=COLOR_ACCENT,  size=12, italic=False)
    add_caption(s, "3 ·  Reset password",  xs[2], Inches(5.0), img_w, color=COLOR_SUCCESS, size=12, italic=False)

    bullets_data = [
        (xs[0], img_w, ["Email + password, remember-me", "Modern glass-card shell", "Logs: login / login_failed"]),
        (xs[1], img_w, ["Admin → reset link (60 min)", "Non-admin → contact-admin panel", "Same response → no email leak"]),
        (xs[2], img_w, ["Strength meter + match indicator", "Requirements checklist", "Re-checks role on submit"]),
    ]
    for x, w, lines in bullets_data:
        box = s.shapes.add_textbox(x, Inches(5.4), w, Inches(1.6))
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = "●  "
            style_run(r, size=8, bold=True, color=COLOR_PRIMARY)
            r2 = p.add_run(); r2.text = line
            style_run(r2, size=10, color=COLOR_TEXT_2)
            p.space_after = Pt(2)


def slide10_profile(prs, n=13):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Self-service profile  ·  My Account")
    add_modern_subtitle(s, "Every authenticated user updates their own basics — no admin needed.")

    add_bullets(s, [
        ("Name & Email",          "instant update, unique enforced"),
        ("Avatar upload",         "JPG / PNG / WebP up to 2 MB"),
        ("Change Password",       "requires current_password"),
        ("Role / Module Access",  "shown read-only"),
        ("Activity log",          "profile_updated with changed_fields"),
    ], left=Inches(0.55), top=Inches(2.5), width=Inches(5.5), size=13, spacing=8)
    add_pill(s, Inches(0.6), Inches(5.6), Inches(5.4), Inches(0.5),
             "Role & permissions  →  admin only", COLOR_FAINT, size=10)

    add_screenshot_rounded(s, "16-profile.png", Inches(6.4), Inches(2.4), Inches(6.4), Inches(4.0))
    add_caption(s, "/profile — My Account page", Inches(6.4), Inches(6.45), Inches(6.4))


def slide11_dashboard(prs, n=15):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Dashboard")
    add_modern_subtitle(s, "Visual digest of the whole estate, refreshed on every load.")

    add_screenshot_rounded(s, "10-dashboard.png", Inches(0.55), Inches(2.4), Inches(8.0), Inches(5.0))
    add_caption(s, "/dashboard — KPIs, inventory chart, recent activity, expiring lists",
                Inches(0.55), Inches(7.45), Inches(8.0))

    add_bullets(s, [
        ("4 KPI cards",        "PCs · Devices · Subs · Expiring"),
        ("Inventory chart",    "Chart.js — PC vs Device by status"),
        ("Recent activity",    "last 8 audit-log entries"),
        ("Expiring tables",    "subs and licenses in next 30 days"),
        ("Quick actions",      "header buttons jump to each module"),
    ], left=Inches(8.85), top=Inches(2.4), width=Inches(4.2), size=11, spacing=7)


def slide12_notifications(prs, n=16):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Notifications system")
    add_modern_subtitle(s, "Live count · per-user read state · auto-unread on signature drift.")

    add_screenshot_rounded(s, "15-notifications.png", Inches(0.55), Inches(2.4), Inches(8.0), Inches(5.0))
    add_caption(s, "/notifications — live items with module colors and per-row actions",
                Inches(0.55), Inches(7.45), Inches(8.0))

    add_bullets(s, [
        ("Topbar bell",         "live count per user, urgency tone"),
        ("Filter chips",        "All · Unread · Read"),
        ("Per-module filter",   "subs · L&C"),
        ("Mark as read",        "single item or all"),
        ("Auto-unread",         "on renewal or bucket transition"),
    ], left=Inches(8.85), top=Inches(2.4), width=Inches(4.2), size=11, spacing=7)


def slide13_mail(prs, n=18):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Mail & staggered reminders")
    add_modern_subtitle(s, "DB-overridable SMTP plus one digest per (module × selected day-mark).")

    img_w = Inches(6.05); img_h = Inches(3.78)
    add_screenshot_rounded(s, "17-mail-settings.png",         Inches(0.55), Inches(2.4), img_w, img_h)
    add_caption(s, "/mail-settings — SMTP, auth, from identity, test send", Inches(0.55), Inches(6.25), img_w)
    add_screenshot_rounded(s, "18-notification-settings.png", Inches(6.75), Inches(2.4), img_w, img_h)
    add_caption(s, "/notification-settings — per-module enable + 30/20/10 day-mark toggles",
                Inches(6.75), Inches(6.25), img_w)


def slide14_audit(prs, n=19):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Activity log  ·  audit trail")
    add_modern_subtitle(s, "Every meaningful action recorded via App\\Support\\ActivityLogger::log(...).")

    add_bullets(s, [
        ("CRUD on every module",           "created · updated · deleted · imported"),
        ("Subscriptions",                  "renewed"),
        ("Auth",                           "login · logout · login_failed · password_reset_*"),
        ("User & profile",                 "profile_updated · mail_sent · mail_test"),
        ("Settings",                       "Mail / Notification updates"),
        ("Captured fields",                "subject_type/id · user_id/name/email · properties · ip · user_agent"),
        ("Visible at",                     "/activity-logs (admin-only, paginated)"),
    ], top=Inches(2.5), size=13, spacing=8)


def slide15_setup(prs, n=21):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Setup & next steps")
    add_modern_subtitle(s, "How to get it running, and where it could go next.")

    add_flat_card(s, Inches(0.55), Inches(2.4), Inches(6.05), Inches(4.6),
                  "Get it running", [
                      "composer install",
                      "cp .env.example .env  →  set APP_*, DB_*, MAIL_*",
                      "php artisan key:generate",
                      "php artisan migrate",
                      "php artisan storage:link",
                      "npm install && npm run build  (optional)",
                      "Schedule:  php artisan schedule:run  (every minute)",
                  ], accent=COLOR_PRIMARY)

    add_flat_card(s, Inches(6.75), Inches(2.4), Inches(6.05), Inches(4.6),
                  "Possible next iterations", [
                      "Warranty-expiry tracking for PCs & Devices",
                      "Two-factor authentication for admins",
                      "Weekly overdue digest",
                      "Per-user notification preferences",
                      "PDF export for contracts",
                      "API surface for integrations",
                      "Soft-delete + restore on critical tables",
                  ], accent=COLOR_ACCENT)


def _member_avatar(slide, x, y, size, initials, gradient_a, gradient_b):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = gradient_a
    circle.line.color.rgb = gradient_b
    circle.line.width = Pt(2.5)
    add_soft_shadow(circle, blur_pt=20, dist_pt=6, color=gradient_a, alpha=25)
    tf = circle.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = initials
    style_run(r, size=38, bold=True, color=COLOR_WHITE)


def slide16_members(prs, n=3):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    apply_template_chrome(s, prs, n)
    add_modern_title(s, "Team")
    add_modern_subtitle(s, "The people who built and maintain this system.")

    members = [
        ("Thiha Zin Ko",            "TZ", COLOR_PRIMARY, COLOR_ACCENT),
        ("Zwe Myat Thurain Phyo",   "ZM", COLOR_ACCENT,  COLOR_DANGER),
        ("Nyein Wai Khant",         "NW", COLOR_SUCCESS, COLOR_PRIMARY),
    ]

    card_w = Inches(3.6); card_h = Inches(4.0)
    gap = Inches(0.4)
    total_w = card_w * 3 + gap * 2
    start_x = (prs.slide_width - total_w) / 2
    y = Inches(2.45)

    for i, (name, initials, c_a, c_b) in enumerate(members):
        x = start_x + i * (card_w + gap)
        # card (no top stripe, soft shadow)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.adjustments[0] = 0.06
        card.fill.solid(); card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_HAIRLINE
        card.line.width = Pt(0.75)
        add_soft_shadow(card, blur_pt=22, dist_pt=8, alpha=15)
        # avatar
        av_size = Inches(1.75)
        av_x = x + (card_w - av_size) / 2
        _member_avatar(s, av_x, y + Inches(0.5), av_size, initials, c_a, c_b)
        # name
        name_box = s.shapes.add_textbox(x + Inches(0.2), y + Inches(2.6), card_w - Inches(0.4), Inches(0.7))
        tf = name_box.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = name
        style_run(r, size=16, bold=True, color=COLOR_TEXT)
        # role
        role_box = s.shapes.add_textbox(x + Inches(0.2), y + Inches(3.3), card_w - Inches(0.4), Inches(0.4))
        rp = role_box.text_frame.paragraphs[0]; rp.alignment = PP_ALIGN.CENTER
        r2 = rp.add_run(); r2.text = "Team member"
        style_run(r2, size=10, color=COLOR_MUTED)


def slide17_thanks(prs):
    s = hero_slide(prs)
    # Big "Thank you"
    title = s.shapes.add_textbox(
        Inches(0.5),
        (prs.slide_height - Inches(3)) / 2,
        prs.slide_width - Inches(1.0),
        Inches(1.7),
    )
    p = title.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Thank you"
    style_run(r, size=88, bold=True, color=COLOR_WHITE)

    # underline-ish accent — saturated golden cream
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              (prs.slide_width - Inches(1.2)) / 2,
                              (prs.slide_height - Inches(3)) / 2 + Inches(1.75),
                              Inches(1.2), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xFC, 0xD3, 0x4D)
    line.line.fill.background()

    # subtitle — pure white
    sub = s.shapes.add_textbox(
        Inches(0.5),
        (prs.slide_height - Inches(3)) / 2 + Inches(2.0),
        prs.slide_width - Inches(1.0),
        Inches(0.7),
    )
    p = sub.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Questions  ·  Discussion"
    style_run(r, size=22, color=COLOR_WHITE)

    # tagline
    tag = s.shapes.add_textbox(
        Inches(0.5),
        prs.slide_height - Inches(1.0),
        prs.slide_width - Inches(1.0),
        Inches(0.5),
    )
    p = tag.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Infra Ninja  ·  IT Assets Management System"
    style_run(r, size=12, color=RGBColor(0xE0, 0xE7, 0xFF))


# ============================================================
def main():
    out_dir = Path(__file__).resolve().parent
    out_path = out_dir / "InfraNinja_ITAMS_Overview.pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Overview chapter (no divider before — title acts as the cover) ----
    slide1_title(prs)                                  # 01
    slide2_overview(prs,      n=2)                     # 02
    slide16_members(prs,      n=3)                     # 03  Team — moved up
    slide3_stack(prs,         n=4)                     # 04
    slide4_arch(prs,          n=5)                     # 05

    # ---- Modules chapter ----
    slide_divider(prs, page_num=6,  chapter_num=1,
                  chapter_label="Modules",
                  description="The four asset categories your team manages — end-to-end CRUD with audit log and bulk operations.",
                  color=COLOR_SUCCESS)
    slide5_modules(prs,       n=7)                     # 07
    slide6_pc_devices(prs,    n=8)                     # 08
    slide7_subs_lic(prs,      n=9)                     # 09

    # ---- Users & Auth chapter ----
    slide_divider(prs, page_num=10, chapter_num=2,
                  chapter_label="Users & Auth",
                  description="Two roles, granular per-module permissions, and a modern login / reset flow.",
                  color=COLOR_PRIMARY)
    slide8_perms(prs,         n=11)                    # 11
    slide9_auth(prs,          n=12)                    # 12
    slide10_profile(prs,      n=13)                    # 13

    # ---- Insights chapter ----
    slide_divider(prs, page_num=14, chapter_num=3,
                  chapter_label="Insights",
                  description="What the system shows at a glance — dashboard, notifications, live counts.",
                  color=COLOR_WARN)
    slide11_dashboard(prs,    n=15)                    # 15
    slide12_notifications(prs, n=16)                   # 16

    # ---- Operations chapter ----
    slide_divider(prs, page_num=17, chapter_num=4,
                  chapter_label="Operations",
                  description="The plumbing keeping it running — mail, reminders, and the audit trail.",
                  color=COLOR_DANGER)
    slide13_mail(prs,         n=18)                    # 18
    slide14_audit(prs,        n=19)                    # 19

    # ---- Wrap-up chapter ----
    slide_divider(prs, page_num=20, chapter_num=5,
                  chapter_label="Wrap-up",
                  description="Setup, where it could go next, and the team behind it.",
                  color=COLOR_MUTED)
    slide15_setup(prs,        n=21)                    # 21
    slide17_thanks(prs)                                # 22

    prs.save(out_path)
    print(f"Wrote {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
